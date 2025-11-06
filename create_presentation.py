from pptx import Presentation
from pptx.util import Inches
import os

os.makedirs('presentacion', exist_ok=True)

prs = Presentation()
# Slide 1: portada
slide = prs.slides.add_slide(prs.slide_layouts[5])
try:
    title = slide.shapes.title
    title.text = "Arcor — Resumen ejecutivo"
except Exception:
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    tx.text_frame.text = "Arcor — Resumen ejecutivo"

# Slide 2: añadir imagen del dashboard
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
try:
    title2 = slide2.shapes.title
    title2.text = "Resumen visual del dashboard"
except Exception:
    pass

img_path = 'resultados/resumen_visual_dashboard.png'
if os.path.exists(img_path):
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    slide2.shapes.add_picture(img_path, left, top, width=width)
else:
    tx2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tx2.text_frame.text = f"Imagen no encontrada: {img_path}. Generá el PNG desde el notebook."

# Guardar presentación
prs.save('presentacion/presentacion_powerpoint.pptx')
print('PPTX guardado en presentacion/presentacion_powerpoint.pptx')
